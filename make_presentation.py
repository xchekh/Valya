from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)   # title / heading bg
MID_BLUE    = RGBColor(0x2E, 0x6D, 0xA4)   # accent
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # slide bg tint
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD = RGBColor(0xE8, 0xA8, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ──────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size=18, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_bullet_box(slide, l, t, w, h, items, title=None,
                   title_size=20, bullet_size=16,
                   title_color=DARK_BLUE, bullet_color=DARK_TEXT,
                   line_spacing=1.15):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = item.get("level", 0)
        run = p.add_run()
        bullet_char = "•  " if p.level == 0 else "   –  "
        run.text = bullet_char + item["text"]
        run.font.size = Pt(item.get("size", bullet_size))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = item.get("color", bullet_color)
        run.font.name = "Calibri"
    return txb


def bg(slide, color=LIGHT_BLUE):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=color)


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, fill_rgb=DARK_BLUE)
    add_textbox(slide, 0.35, 0.12, 12.6, 0.7, title,
                font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.35, 0.82, 12.6, 0.45, subtitle,
                    font_size=18, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)


def footer(slide, text="CS 435 · Newton's Method · Nisheeth Vishnoi (2018)"):
    add_rect(slide, 0, 7.18, 13.33, 0.32, fill_rgb=DARK_BLUE)
    add_textbox(slide, 0.3, 7.19, 12.7, 0.28, text,
                font_size=11, color=WHITE, align=PP_ALIGN.LEFT)


def info_box(slide, l, t, w, h, title, body_lines,
             bg_color=WHITE, title_color=MID_BLUE,
             body_color=DARK_TEXT, title_size=17, body_size=15):
    add_rect(slide, l, t, w, h, fill_rgb=bg_color,
             line_rgb=MID_BLUE, line_width=1.2)
    txb = slide.shapes.add_textbox(
        Inches(l+0.12), Inches(t+0.1),
        Inches(w-0.24), Inches(h-0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Calibri"
    for line in body_lines:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = line
        run2.font.size = Pt(body_size)
        run2.font.color.rgb = body_color
        run2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, DARK_BLUE)
add_rect(sl, 0.6, 1.2, 12.13, 4.2, fill_rgb=RGBColor(0x0D,0x2A,0x45),
         line_rgb=ACCENT_GOLD, line_width=2)
add_textbox(sl, 0.9, 1.45, 11.5, 1.1,
            "Newton's Method", font_size=54, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.9, 2.6, 11.5, 0.55,
            "A Second-Order Optimization Algorithm",
            font_size=26, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.9, 3.3, 11.5, 0.45,
            "Derivation · Convergence Analysis · Riemannian Perspective",
            font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 0.9, 4.05, 11.5, 0.4,
            "Based on: CS 435 Lecture 6  ·  Nisheeth Vishnoi  ·  April 12, 2018",
            font_size=15, color=RGBColor(0xAA,0xCC,0xEE), align=PP_ALIGN.CENTER)
add_rect(sl, 0, 7.18, 13.33, 0.32, fill_rgb=RGBColor(0x0A,0x1E,0x35))
add_textbox(sl, 0.3, 7.19, 12.7, 0.28,
            "Scientific Seminar Presentation",
            font_size=12, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – Outline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Lecture Outline", "Five major topics covered")
footer(sl)

topics = [
    ("1", "Newton-Raphson Method",          "Root finding, update rule, quadratic convergence, multivariate extension"),
    ("2", "Unconstrained Optimization",     "From root-finding to minimization; second-order approximation viewpoint"),
    ("3", "Euclidean-Norm Analysis",        "Theorem 2, condition NE, limitations & affine invariance issue"),
    ("4", "Riemannian Manifold View",       "Hessian manifold, local norm, Newton's flow as gradient descent"),
    ("5", "Local-Norm Convergence",         "Condition NL, Theorem 4: affinely invariant quadratic convergence"),
]
x0, y0, bw, bh = 0.35, 1.5, 12.63, 0.95
for i,(num,title,desc) in enumerate(topics):
    yy = y0 + i*(bh+0.06)
    add_rect(sl, x0, yy, bw, bh, fill_rgb=WHITE,
             line_rgb=MID_BLUE, line_width=1)
    add_rect(sl, x0, yy, 0.65, bh, fill_rgb=MID_BLUE)
    add_textbox(sl, x0+0.05, yy+0.18, 0.55, 0.55, num,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, x0+0.75, yy+0.05, 4.5, 0.45, title,
                font_size=18, bold=True, color=DARK_BLUE)
    add_textbox(sl, x0+0.75, yy+0.5, 11.5, 0.4, desc,
                font_size=14, color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – Newton-Raphson: Idea & Update Rule
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Newton-Raphson Method", "Root finding via tangent-line linearisation")
footer(sl)

add_bullet_box(sl, 0.35, 1.5, 6.0, 3.2,
    title="Core Idea",
    items=[
        {"text": "Given g : ℝ → ℝ, find r such that g(r) = 0"},
        {"text": "Start at x₀ near a root; draw the tangent line at (x₀, g(x₀))"},
        {"text": "Let x₁ be where that tangent meets the x-axis"},
        {"text": "Repeat: quadratic approach to the root"},
    ],
    title_size=19, bullet_size=16)

info_box(sl, 0.35, 4.85, 6.0, 1.3,
         "Update Rule",
         ["x_{k+1} := x_k  −  g(x_k) / g′(x_k)    for all k ≥ 0",
          "Requires: g twice continuously differentiable (C²)"],
         bg_color=RGBColor(0xE8,0xF4,0xFF), title_size=17, body_size=15)

info_box(sl, 6.7, 1.5, 6.28, 4.65,
         "Worked Example  (computing 1/a)",
         ["Minimise f(x) = ax − log x  ⟹  g(x) = a − 1/x",
          "",
          "Newton update: x_{k+1} = 2x_k − a·x_k²",
          "  (no division required — only + and ×!)",
          "",
          "Let e_k = 1 − a·x_k.  Then  e_{k+1} = e_k²",
          "",
          "Converges iff |e₀| < 1  (i.e. 0 < x₀ < 2/a)",
          "Key insight: initial point choice is CRITICAL"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – Quadratic Convergence
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Quadratic Convergence (Theorem 1)",
           "Distance to root is squared at every iteration")
footer(sl)

info_box(sl, 0.35, 1.5, 12.63, 1.55,
         "Theorem 1  (Quadratic Convergence for Root Finding)",
         ["Let g : ℝ → ℝ be C², r a root, x₀ a starting point, x₁ = x₀ − g(x₀)/g′(x₀).  Then:",
          "   |r − x₁|  ≤  M · |r − x₀|²      where   M = sup_{ξ ∈ [r,x₀]}  |g″(ξ) / (2g′(x₀))|"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=18, body_size=16)

add_bullet_box(sl, 0.35, 3.2, 6.1, 2.0,
    title="Proof Sketch (Taylor expansion)",
    items=[
        {"text": "Expand g(r) = 0 around x₀ via Taylor / MVT"},
        {"text": "Use g(x₀) = g′(x₀)(x₀ − x₁) from the update definition"},
        {"text": "Algebra yields: g′(x₀)(r−x₁) = ½(r−x₀)²g″(ξ)"},
        {"text": "Divide by |g′(x₀)| to get the bound"},
    ],
    title_size=18, bullet_size=15)

info_box(sl, 6.7, 3.2, 6.28, 2.0,
         "Convergence Rate",
         ["If M ≤ 1 and |x₀ − r| < ½ :",
          "",
          "   |x_k − r|  ≤  |x₀ − r|^{2^k}  ≤  2^{−2^k}",
          "",
          "To reach error ε: only k ≈ log log(1/ε) steps!",
          "  ⟹  doubly-exponential speed"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=15)

add_textbox(sl, 0.35, 5.35, 12.63, 0.55,
            "Practical note: Newton's method is robust even when no explicit bounds on M or |x₀ − r| are available.",
            font_size=14, color=MID_BLUE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – Multivariate Extension
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Extension to Multivariate Functions",
           "Systems of nonlinear equations in ℝⁿ")
footer(sl)

add_bullet_box(sl, 0.35, 1.5, 6.1, 2.1,
    title="Setting",
    items=[
        {"text": "Find x ∈ ℝⁿ satisfying g(x) = 0,  g : ℝⁿ → ℝⁿ"},
        {"text": "g(x) = (g₁(x), …, gₙ(x))ᵀ  —  n equations, n unknowns"},
        {"text": "Jacobian Jg(x₀) replaces the scalar derivative g′(x₀)"},
        {"text": "g(x) ≈ g(x₀) + Jg(x₀)(x − x₀) + o(‖x − x₀‖²)"},
    ],
    title_size=18, bullet_size=15)

info_box(sl, 0.35, 3.75, 6.1, 1.45,
         "Multivariate Update Rule",
         ["x_{k+1} := x_k − Jg(x_k)⁻¹ · g(x_k)     for all k ≥ 0",
          "",
          "Local quadratic convergence still holds"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

info_box(sl, 6.7, 1.5, 6.28, 3.7,
         "Jacobian vs. Scalar Derivative",
         ["Univariate:  x₁ = x₀ − g(x₀) / g′(x₀)",
          "",
          "Multivariate:",
          "  • g(x₀) is a vector  ∈ ℝⁿ",
          "  • Jg(x₀) is an n×n matrix  [∂gᵢ/∂xⱼ]",
          "  • Division becomes matrix inversion",
          "",
          "Convergence rate M now involves:",
          "  • Upper bound on Lipschitz const. of Jg",
          "  • Lower bound on ‖Jg(x)⁻¹‖ (spectral norm)"],
         bg_color=WHITE, title_size=17, body_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – Newton's Method for Optimization
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Newton's Method for Unconstrained Optimization",
           "From root-finding to minimisation via gradient zeroing")
footer(sl)

info_box(sl, 0.35, 1.5, 12.63, 1.2,
         "Key Bridge",
         ["Minimising a smooth convex f  ≡  finding a root of ∇f",
          "because  x* = argmin f(x)  iff  ∇f(x*) = 0  (for convex f)"],
         bg_color=RGBColor(0xE8,0xF4,0xFF), title_size=17, body_size=16)

add_bullet_box(sl, 0.35, 2.85, 5.9, 2.5,
    title="Newton Step & Update",
    items=[
        {"text": "Apply multivariate rule to g := ∇f"},
        {"text": "Jacobian of ∇f  =  Hessian ∇²f"},
        {"text": "Update:  x_{k+1} := x_k − (∇²f(x_k))⁻¹ ∇f(x_k)"},
        {"text": "Define Newton step:  n(x) := −(∇²f(x))⁻¹ ∇f(x)"},
        {"text": "Compact form:  x_{k+1} = x_k + n(x_k)"},
    ],
    title_size=18, bullet_size=15)

info_box(sl, 6.7, 2.85, 6.28, 2.5,
         "Second-Order Interpretation",
         ["At each step, minimise the quadratic approximation:",
          "",
          "  f̃(x) = f(x₀) + ⟨x−x₀, ∇f(x₀)⟩ + ½(x−x₀)ᵀ∇²f(x₀)(x−x₀)",
          "",
          "Setting ∇f̃(x) = 0 yields exactly the Newton step!",
          "",
          "⟹ One step suffices for strictly convex quadratics"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=14)

add_bullet_box(sl, 0.35, 5.5, 12.63, 1.0,
    title="Computational Cost per Iteration",
    items=[
        {"text": "Must solve ∇²f(xₖ) · d = ∇f(xₖ)  —  an n×n linear system"},
        {"text": "Gaussian elimination: O(n³);  fast MM: O(n^ω);  special structures (SDD, Laplacians): Õ(m)  [Spielman–Teng 2004]"},
    ],
    title_size=17, bullet_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – Logarithmic Barrier & Analytic Center
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Example: Logarithmic Barrier & Analytic Center",
           "Newton's method applied to a polytope interior-point function")
footer(sl)

info_box(sl, 0.35, 1.5, 6.1, 1.5,
         "Polytope Setting",
         ["P = { x ∈ ℝⁿ : ⟨aᵢ, x⟩ ≤ bᵢ,  i = 1,…,m }",
          "Bounded, full-dimensional subset of ℝⁿ"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

info_box(sl, 6.7, 1.5, 6.28, 1.5,
         "Logarithmic Barrier",
         ["F(x) = − Σᵢ log(bᵢ − aᵢᵀ x)",
          "Each term → +∞ as x approaches face { aᵢᵀx = bᵢ }"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=15)

info_box(sl, 0.35, 3.15, 6.1, 2.1,
         "Gradient & Hessian",
         ["sᵢ(x) := bᵢ − aᵢᵀx   (slack for constraint i)",
          "",
          "∇F(x) = Σᵢ  aᵢ / sᵢ(x)",
          "",
          "∇²F(x) = Σᵢ  aᵢaᵢᵀ / sᵢ(x)²"],
         bg_color=WHITE, title_size=17, body_size=15)

info_box(sl, 6.7, 3.15, 6.28, 2.1,
         "Analytic Center",
         ["x* with ∇F(x*) = 0  is the analytic center of P",
          "",
          "Physical view: point where all constraint 'forces' balance",
          "",
          "Used to: assess polytope volume, initialise LP solvers",
          "⟹ Key primitive for polynomial-time Linear Programming"],
         bg_color=WHITE, title_size=17, body_size=14)

add_bullet_box(sl, 0.35, 5.4, 12.63, 0.85,
    title="Computational Complexity",
    items=[
        {"text": "∇F(x): O(nm);  ∇²F(x): O(n²m) — reduces to O(m) for sparse aᵢ"},
        {"text": "If aᵢ form a graph incidence matrix ⟹ ∇²F is a Laplacian, solvable in Õ(m) [Spielman–Teng]"},
    ],
    title_size=17, bullet_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – Euclidean Norm Analysis
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Convergence Analysis: Euclidean Norm",
           "Theorem 2 and Condition NE")
footer(sl)

info_box(sl, 0.35, 1.5, 12.63, 1.5,
         "Theorem 2  (Quadratic Convergence w.r.t. Euclidean Norm)",
         ["Let f : ℝⁿ → ℝ, x* its minimiser, x₀ any starting point, x₁ := x₀ + n(x₀).",
          "If condition NE(M) holds then:  ‖x₁ − x*‖₂  ≤  M · ‖x₀ − x*‖₂²"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

info_box(sl, 0.35, 3.15, 7.1, 2.35,
         "Condition NE(M) — Newton-Euclidean",
         ["∃ ball B(x*, R) containing x₀,  constants h, L > 0  with  M ≥ L/(2h)  s.t.:",
          "",
          "  1. ‖H(x)⁻¹‖ ≤ 1/h   for all x ∈ B(x*, R)",
          "     (Hessian 'magnitude' is large enough)",
          "",
          "  2. ‖H(x) − H(y)‖ ≤ L‖x−y‖₂   for all x,y ∈ B(x*, R)",
          "     (Hessian is Lipschitz-continuous)"],
         bg_color=WHITE, title_size=17, body_size=14)

info_box(sl, 7.65, 3.15, 5.33, 2.35,
         "Proof Strategy",
         ["Apply FTC to ∇f to expand x₁ − x*:",
          "",
          "x₁ − x* = H(x₀)⁻¹ ∫₀¹ (H(x₀+t(x*−x₀)) − H(x₀))(x*−x₀) dt",
          "",
          "Take norms; use Lipschitz condition on integral.",
          "Result: ‖x₁−x*‖₂ ≤ (L‖H(x₀)⁻¹‖ / 2) ‖x*−x₀‖₂²"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=13)

add_bullet_box(sl, 0.35, 5.65, 12.63, 0.85,
    title="Analogy with Theorem 1",
    items=[
        {"text": "Condition 1 (large h): Hessian plays role of |g′(x)| being bounded away from zero"},
        {"text": "Condition 2 (Lipschitz H): bounds magnitude of third-order term D³f"},
    ],
    title_size=17, bullet_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 – Limitations & Affine Invariance
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Limitations of Euclidean-Norm Analysis",
           "Affine invariance and why the Euclidean norm is the wrong choice")
footer(sl)

info_box(sl, 0.35, 1.5, 7.1, 2.7,
         "Failure Example: Elongated Polytope",
         ["P = [−K₁, K₁] × [−1/K₂, 1/K₂] ⊆ ℝ²,   K₁, K₂ ≫ 1",
          "",
          "Hessian at x* = (0,0):  h ≤ 2/K₁²  (tiny!)",
          "",
          "Lipschitz constant of H:  L ≈ K₂²  (huge!)",
          "",
          "M = L/(2h) = Ω(K₁²K₂²)  — Theorem 2 gives no useful bound",
          "",
          "Yet Newton's method converges rapidly in practice!"],
         bg_color=RGBColor(0xFF,0xEB,0xE6), title_size=17, body_size=14)

info_box(sl, 7.65, 1.5, 5.33, 2.7,
         "Affine Invariance of Newton's Method",
         ["Consider affine map  φ(x) = Ax + b",
          "",
          "If x₀ → x₁ under Newton's method for f̃(x) = f(Ax+b)",
          "then φ(x₀) → φ(x₁) under Newton's method for f",
          "",
          "⟹ Newton's trajectory is coordinate-independent!",
          "",
          "Gradient descent does NOT have this property",
          "(preconditioning can improve GD convergence)"],
         bg_color=WHITE, title_size=17, body_size=14)

info_box(sl, 0.35, 4.35, 12.63, 1.35,
         "The Core Problem",
         ["The constants h and L in condition NE change under affine reparametrisation,",
          "even though Newton's method takes the same steps.",
          "",
          "Solution: switch to an affinely invariant norm — the LOCAL NORM induced by the Hessian."],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 – Riemannian Manifold / Local Norm
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Newton's Method as Gradient Descent on a Riemannian Manifold",
           "The Hessian manifold and the local norm")
footer(sl)

add_bullet_box(sl, 0.35, 1.5, 6.1, 2.3,
    title="Motivating Observation",
    items=[
        {"text": "Minimise f(x₁,x₂) = x₁² + Kx₂²  (K large)"},
        {"text": "Gradient descent with fixed η: overshoots (‖∇f‖ ≈ 2K)"},
        {"text": "Must use  η ≈ 1/K  — very slow convergence"},
        {"text": "Better idea: replace Euclidean norm with the norm"},
        {"text": "   ‖(u₁,u₂)‖◦ = √(u₁²+Ku₂²) (the Hessian norm)"},
        {"text": "With this norm, the steepest descent points directly to x*"},
        {"text": "For η=1: one step reaches the exact minimum!"},
    ],
    title_size=18, bullet_size=14)

info_box(sl, 6.7, 1.5, 6.28, 2.3,
         "Local Inner Product & Local Norm",
         ["For strictly convex f with Hessian H(x) = ∇²f(x):",
          "",
          "  ⟨u, v⟩_x := uᵀH(x)v",
          "  ‖u‖_x := √(uᵀH(x)u)",
          "",
          "These define a Riemannian metric on ℝⁿ",
          "that varies pointwise with x",
          "",
          "Affinely invariant  ✓"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

info_box(sl, 0.35, 3.95, 12.63, 1.6,
         "Newton's Flow = Gradient Descent on Hessian Manifold",
         ["Steepest descent w.r.t. local norm ‖·‖_x :",
          "   max_{‖u‖_x ≤ 1} [−⟨∇f(x), u⟩]   ⟹   u_opt ∝ −H(x)⁻¹∇f(x)  =  n(x)",
          "",
          "Continuous-time limit:  dx/dt = −(∇²f(x))⁻¹ ∇f(x)   ← Newton's flow",
          "",
          "Newton's method = Euler discretisation of Newton's flow on a Hessian manifold"],
         bg_color=WHITE, title_size=17, body_size=15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 – Local Norm Convergence
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Affinely Invariant Convergence: Local Norm Analysis",
           "Theorem 4 and Condition NL")
footer(sl)

info_box(sl, 0.35, 1.5, 12.63, 1.15,
         "New Potential Function  λ(x) := ‖n(x)‖_x",
         ["λ(x) = √(∇f(x)ᵀ H(x)⁻¹ ∇f(x))       — affinely invariant",
          "Interpretation: ½λ(x)² = gap between f(x) and minimum of the local quadratic approximation"],
         bg_color=RGBColor(0xE2,0xF0,0xFF), title_size=17, body_size=15)

info_box(sl, 0.35, 2.8, 12.63, 1.35,
         "Theorem 4  (Quadratic Convergence w.r.t. Local Norm)",
         ["Let f satisfy condition NL, x₀ ∈ ℝⁿ, x₁ := x₀ + n(x₀).",
          "",
          "If  λ(x₀) = ‖n(x₀)‖_{x₀} < 1/6   then:",
          "   ‖n(x₁)‖_{x₁}  ≤  3 · ‖n(x₀)‖_{x₀}²"],
         bg_color=RGBColor(0xE8,0xF4,0xFF), title_size=17, body_size=16)

info_box(sl, 0.35, 4.3, 6.1, 1.85,
         "Condition NL (Newton-Local)  — affinely invariant",
         ["For all x, y with  ‖y − x‖_x = δ < 1 :",
          "",
          "  (1 − 3δ) H(x)  ⪯  H(y)  ⪯  (1 + 3δ) H(x)",
          "",
          "Hessian changes at most proportionally to itself",
          "when moving by one 'local unit'"],
         bg_color=WHITE, title_size=17, body_size=14)

info_box(sl, 6.7, 4.3, 6.28, 1.85,
         "Proof Key Steps",
         ["1. Lemma 6: nearby local norms differ by ≤ 2×",
          "   (‖y−x‖_x ≤ 1/6 ⟹ ½‖·‖_x ≤ ‖·‖_y ≤ 2‖·‖_x)",
          "",
          "2. Write ∇f(x₁) = M(x₀) H(x₀)⁻¹ ∇f(x₀)  via FTC",
          "",
          "3. Use NL + Fact 1 to bound ‖H⁻¹/²M H⁻¹/²‖",
          "   ≤ (3/2)λ(x₀)  ⟹  Theorem 4"],
         bg_color=RGBColor(0xFF,0xFB,0xE6), title_size=17, body_size=14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 – Summary & Comparison
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Summary & Comparison of Analyses",
           "Two convergence frameworks side by side")
footer(sl)

# Table header
add_rect(sl, 0.35, 1.5, 12.63, 0.45, fill_rgb=DARK_BLUE)
for txt, x in [("Property", 0.4), ("Theorem 2 (Euclidean)", 3.5), ("Theorem 4 (Local Norm)", 8.1)]:
    add_textbox(sl, x, 1.53, 4.5, 0.38, txt,
                font_size=15, bold=True, color=WHITE)

rows = [
    ("Convergence type",        "‖x₁ − x*‖₂ ≤ M‖x₀−x*‖₂²",      "‖n(x₁)‖_{x₁} ≤ 3‖n(x₀)‖²_{x₀}"),
    ("Convergence condition",   "Condition NE (h, L, ball radius)", "Condition NL (local Lipschitz)"),
    ("Affinely invariant?",     "No  ✗  (h, L change with coords)", "Yes  ✓  (all quantities invariant)"),
    ("Fails for log-barrier?",  "Yes — M = Ω(K₁²K₂²) unusable",   "No  ✓  — bounded regardless of P"),
    ("Proof tool",              "Lipschitz H + operator norms",     "FTC + Lemma 6 + PSD ordering"),
]
for i, (prop, col2, col3) in enumerate(rows):
    yy = 2.0 + i*0.88
    bg_c = WHITE if i % 2 == 0 else RGBColor(0xF0,0xF7,0xFF)
    add_rect(sl, 0.35, yy, 12.63, 0.88, fill_rgb=bg_c,
             line_rgb=RGBColor(0xCC,0xDD,0xEE), line_width=0.5)
    add_textbox(sl, 0.4, yy+0.08, 3.0, 0.72, prop,
                font_size=13, bold=True, color=DARK_BLUE)
    add_textbox(sl, 3.5, yy+0.08, 4.5, 0.72, col2,
                font_size=13, color=DARK_TEXT)
    add_textbox(sl, 8.1, yy+0.08, 4.7, 0.72, col3,
                font_size=13, color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 – Big Picture & Takeaways
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "Big Picture & Key Takeaways", "")
footer(sl)

takeaways = [
    ("Second-Order Power",
     "Newton's method uses both gradient and curvature (Hessian). One step suffices for quadratics; convergence is doubly-exponential near the root."),
    ("Quadratic Convergence",
     "k ≈ log log(1/ε) iterations to reach precision ε — vastly faster than first-order methods which need O(1/ε) or O(log 1/ε) steps."),
    ("Cost vs. Benefit",
     "Each iteration requires solving an n×n linear system (O(n³) in general). Power comes with computational price; special structure (Laplacians, SDD) mitigates this."),
    ("Affine Invariance",
     "Newton's method is coordinate-free. The Euclidean-norm analysis breaks this symmetry; local-norm (Riemannian) analysis restores it via condition NL."),
    ("Connection to Interior-Point Methods",
     "The logarithmic barrier + Newton's method → polynomial-time LP algorithm (Renegar 1988). Condition NL ≈ self-concordance (Nesterov–Nemirovskii 1994)."),
]

for i, (title, body) in enumerate(takeaways):
    y = 1.5 + i * 1.05
    add_rect(sl, 0.35, y, 0.08, 0.75, fill_rgb=ACCENT_GOLD)
    add_rect(sl, 0.55, y, 12.43, 0.95, fill_rgb=WHITE,
             line_rgb=RGBColor(0xCC,0xDD,0xEE), line_width=0.7)
    add_textbox(sl, 0.65, y+0.04, 3.5, 0.4, title,
                font_size=16, bold=True, color=DARK_BLUE)
    add_textbox(sl, 0.65, y+0.46, 12.1, 0.44, body,
                font_size=13, color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 – References
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header_bar(sl, "References", "")
footer(sl)

refs = [
    "[1] A. Galntai. The theory of Newton's method. J. Comput. Appl. Math., 124(1):25–44, 2000.",
    "[2] Y. Nesterov & A. Nemirovskii. Interior-Point Polynomial Algorithms in Convex Programming. SIAM, 1994.",
    "[3] J. Raphson. Analysis aequationum universalis. London, 1690.",
    "[4] J. Renegar. A polynomial-time algorithm, based on Newton's method, for linear programming. Math. Program., 40:59–93, 1988.",
    "[5] D. A. Spielman & S.-H. Teng. Nearly-linear time algorithms for graph partitioning, graph sparsification, and solving linear systems. STOC 2004.",
]
for i, ref in enumerate(refs):
    add_textbox(sl, 0.7, 1.6 + i*0.9, 12.0, 0.75, ref,
                font_size=15, color=DARK_TEXT)

add_textbox(sl, 0.35, 6.4, 12.63, 0.5,
            "Source: CS 435 (2018), Yale University — Lecture 6 notes by Nisheeth Vishnoi",
            font_size=14, color=MID_BLUE, align=PP_ALIGN.CENTER)

# ── save ────────────────────────────────────────────────────────
out = "/home/sony/repo/Valya/Valya/newton_method_presentation.pptx"
prs.save(out)
print("Saved:", out)
